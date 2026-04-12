#!/usr/bin/env python3
"""
generate-synapse-pptx.py — Synapse 分享材料 PPT 生成器
品牌色：Gold #FCAD2A / Deep Blue #013A7D / Cyan #028CDC
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ── Brand Colors ──
GOLD = RGBColor(0xFC, 0xAD, 0x2A)
BLUE = RGBColor(0x01, 0x3A, 0x7D)
CYAN = RGBColor(0x02, 0x8C, 0xDC)
DARK = RGBColor(0x0A, 0x16, 0x28)
CARD = RGBColor(0x11, 0x1D, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0xCB, 0xD5, 0xE1)
MUTED = RGBColor(0x8E, 0x95, 0xA7)
DIM = RGBColor(0x5A, 0x60, 0x72)
RED = RGBColor(0xF5, 0x22, 0x2D)
GREEN = RGBColor(0x52, 0xC4, 0x1A)


def set_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=13, color=TEXT, bold=False, italic=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_multiline(slide, left, top, width, height, lines, font_size=13, color=TEXT, line_spacing=1.5):
    """lines = list of (text, color, bold, size) tuples"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, c, b, sz) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz or font_size)
        p.font.color.rgb = c or color
        p.font.bold = b
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * 0.3)
    return tf


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_gold_bar(slide):
    add_rect(slide, 0, 0, 10, 0.035, GOLD)


def add_label(slide, text):
    add_text(slide, 0.6, 0.15, 5, 0.3, text, font_size=10, color=GOLD, bold=True)


def add_title(slide, text, y=0.5, size=26):
    add_text(slide, 0.6, y, 8.8, 0.7, text, font_size=size, color=WHITE, bold=True)


def add_chat(slide, y, speaker, text, is_user=False):
    border_color = GOLD if is_user else RGBColor(0x1E, 0x29, 0x3B)
    bg_color = RGBColor(0x1A, 0x15, 0x08) if is_user else CARD
    label_color = GOLD if is_user else CYAN
    x = 1.8 if is_user else 0.6
    w = 7.6 if is_user else 7.6

    # Top border
    add_rect(slide, x, y, w, 0.03, border_color)
    # Background
    add_rect(slide, x, y + 0.03, w, 0.75, bg_color)
    # Speaker label
    add_text(slide, x + 0.12, y + 0.05, 1.5, 0.25, speaker, font_size=10, color=label_color, bold=True)
    # Message text
    add_text(slide, x + 0.12, y + 0.25, w - 0.25, 0.5, text, font_size=12, color=TEXT)


def add_annotation(slide, y, text):
    add_rect(slide, 0.6, y, 8.8, 0.42, RGBColor(0x0D, 0x1B, 0x2A))
    add_text(slide, 0.8, y + 0.02, 8.4, 0.38, "💡 " + text, font_size=11, color=GOLD)


def add_expert(slide, y, role, text):
    add_rect(slide, 0.6, y, 1.5, 0.4, CARD)
    add_text(slide, 0.6, y, 1.5, 0.4, role, font_size=10, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, 2.3, y, 7.1, 0.4, text, font_size=11, color=TEXT)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # ═══════ SLIDE 1: 封面 ═══════
    s = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_bg(s)
    add_rect(s, 0, 5.3, 10, 0.04, GOLD)
    add_text(s, 0.6, 0.8, 8.8, 0.9, "SYNAPSE", font_size=52, color=GOLD, bold=True)
    add_text(s, 0.6, 1.8, 8.8, 0.6, "用一场对话，构建AI协作运营体系", font_size=24, color=WHITE, bold=True)
    add_text(s, 0.6, 2.6, 8.8, 0.8, "一个真实案例：如何通过与 AI CEO 的对话，一天内搭建 10 个团队、46 个 Agent 的完整运营体系，\n并在对话过程中不断发现问题、修正机制、持续进化", font_size=13, color=MUTED)
    add_rect(s, 0.6, 3.5, 2, 0.015, GOLD)
    add_text(s, 0.6, 3.7, 5, 0.3, "{{PRESIDENT_NAME}}  |  {{COMPANY_NAME}}  |  2026.04.10", font_size=12, color=DIM)

    # ═══════ SLIDE 2: 背景 ═══════
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_gold_bar(s)
    add_label(s, "背景")
    add_title(s, "AI 技术发展太快，CEO 一个人管不过来")
    add_text(s, 0.6, 1.2, 8.8, 0.4, "{{COMPANY_NAME}} 是一家建筑数字化资产运营管理公司。我们面临所有管理者都面临的挑战：", font_size=13, color=TEXT)
    add_rect(s, 0.6, 1.7, 0.04, 1.1, GOLD)
    add_rect(s, 0.64, 1.7, 8.76, 1.1, CARD)
    add_text(s, 0.85, 1.75, 8.3, 1.0,
        "「我没有足够的时间逐一评估每个AI新工具。\n  我需要一个AI团队替我跟踪技术动态、做评估、做决策、做执行。\n  我只想提出目标，然后验收结果。」",
        font_size=14, color=WHITE, italic=True)
    add_text(s, 0.6, 3.0, 8.8, 0.4, "于是，在一场持续数小时的 Claude Code 对话中，我与 AI CEO Lysander 一起，从零构建了一套完整的 AI 协作运营体系。", font_size=13, color=TEXT)
    add_text(s, 0.6, 3.6, 8.8, 0.3, "以下是这场对话的真实过程 — 包括成功、失败、纠错和进化。", font_size=14, color=GOLD, bold=True)

    # ═══════ SLIDE 3: 场景一 ═══════
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_gold_bar(s)
    add_label(s, "场景一 · 角色定位")
    add_title(s, "第一句话就暴露了问题")
    add_text(s, 0.6, 1.15, 8.8, 0.35, "我要求AI确认管理体系是否生效时，发现系统把总裁和AI CEO的角色混在了一起。", font_size=12, color=MUTED)
    add_chat(s, 1.6, "总裁 {{PRESIDENT_NAME}}", "我是总裁 {{PRESIDENT_NAME}}；Lysander 是CEO，把这个定位确认清楚", is_user=True)
    add_chat(s, 2.55, "Lysander AI", "明白。角色关系已全局更新：总裁{{PRESIDENT_NAME}}=最高决策者，Lysander=AI管理者。已修改5个核心配置文件。", is_user=False)
    add_annotation(s, 3.5, "启示：AI体系的第一步不是技术，是定义谁对谁负责。角色不清，后续所有流程都是空中楼阁。")
    add_text(s, 0.6, 4.1, 8.8, 0.4, "这看似简单的一步，实际上定义了整个体系的权责基础 — 总裁是老板，Lysander是执行CEO，下面有智囊团和执行团队。", font_size=11, color=DIM)

    # ═══════ SLIDE 4: 场景二 — 评审会实录 ═══════
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_gold_bar(s)
    add_label(s, "场景二 · 智囊团评审会实录")
    add_title(s, "总裁说：我参与越少越好", size=24)
    add_chat(s, 1.0, "总裁 {{PRESIDENT_NAME}}", "我参与越少越好，我的判断不一定准确。更多时候应该引入专家评审。我只关注提出目标和最终验收成果。", is_user=True)
    add_text(s, 0.6, 1.95, 8.8, 0.3, "Lysander 立即召集智囊团 4 位专家开评审会，对当前执行链和决策链做完整分析 ——", font_size=12, color=MUTED)
    add_expert(s, 2.35, "战略分析师", "执行链有3个结构性问题：无程序化强制、跨会话状态丢失、简单任务开销过大")
    add_expert(s, 2.8, "决策顾问", "设计完整性 8/10，但强制执行力仅 3/10，跨会话连续性 2/10 — 差距巨大")
    add_expert(s, 3.25, "趋势洞察师", "执行链遵守度会随对话深度递减 — 不是AI偷懒，是架构没有强制检查点")
    add_expert(s, 3.7, "决策顾问", "能否保证每次都按流程执行？诚实回答：以当前架构，不能100%保证。")
    add_annotation(s, 4.25, "评审结果：均分 5.0 全票通过三级改进方案（规则层+机制层+代码层），当天全部落地。")

    # ═══════ SLIDE 5: 场景二续 — 四级决策 ═══════
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_gold_bar(s)
    add_label(s, "场景二 · 决策体系的诞生与纠错")
    add_title(s, "四级决策体系 — 专家先于CEO决策", size=24)
    add_text(s, 0.6, 1.1, 8.8, 0.35, "体系建好后，总裁又发现了一个关键逻辑问题并即时纠正：", font_size=12, color=MUTED)
    add_chat(s, 1.5, "总裁 {{PRESIDENT_NAME}}", "当前的四级决策对吗？我认为 L1自动→L2专家→L3 CEO→L4总裁 更合适", is_user=True)
    # Before/After
    add_rect(s, 0.6, 2.5, 4.2, 1.3, CARD)
    add_rect(s, 0.6, 2.5, 0.04, 1.3, RED)
    add_text(s, 0.8, 2.5, 3, 0.3, "修正前", font_size=12, color=RED, bold=True)
    add_text(s, 0.8, 2.8, 3.8, 0.9, "L1自动 → L2 CEO判断 → L3专家 → L4总裁\n\n问题：CEO没有专家分析就做判断\n可能拍脑袋决策", font_size=11, color=TEXT)
    add_rect(s, 5.2, 2.5, 4.2, 1.3, CARD)
    add_rect(s, 5.2, 2.5, 0.04, 1.3, GREEN)
    add_text(s, 5.4, 2.5, 3, 0.3, "修正后", font_size=12, color=GREEN, bold=True)
    add_text(s, 5.4, 2.8, 3.8, 0.9, "L1自动 → L2专家评审 → L3 CEO决策 → L4总裁\n\nCEO每个决策都有专家分析做支撑\n专业的事先过专家", font_size=11, color=TEXT)
    add_annotation(s, 4.05, "体系设计不是一次成型的。总裁的每次纠错都是体系进化的驱动力。")

    # ═══════ SLIDE 6: 场景三 — 情报系统 ═══════
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_gold_bar(s)
    add_label(s, "场景三 · 从一个需求到全自动管道")
    add_title(s, "每天给我一份有用的AI情报", size=24)
    add_chat(s, 1.0, "总裁 {{PRESIDENT_NAME}}", "AI技术发展太快，我没时间看。目标：每天收到专业AI专家结合我日常工作整理的分析报告，经智囊团审查后通知我。", is_user=True)
    add_text(s, 0.6, 1.95, 8.8, 0.3, "Lysander 接到目标后的执行过程：", font_size=12, color=MUTED)
    steps = [
        ("1.", "新建 AI技术研究员（Graphify第6名成员），定义能力标准和信息源优先级"),
        ("2.", "搜索 5 组 AI 前沿动态（Claude更新/Agent框架/Harness Engineering/工具链/企业案例）"),
        ("3.", "手动生成首期情报日报，开发品牌化HTML报告生成器"),
        ("4.", "创建远程定时任务（每天your local time自动运行）+Slack通知"),
    ]
    for i, (num, text) in enumerate(steps):
        y = 2.35 + i * 0.38
        add_text(s, 0.6, y, 0.3, 0.3, num, font_size=12, color=GOLD, bold=True)
        add_text(s, 0.95, y, 8.45, 0.3, text, font_size=12, color=TEXT)
    add_chat(s, 3.95, "总裁 {{PRESIDENT_NAME}}", "查看了首期报告，质量完全符合我的要求，非常棒。", is_user=True)
    add_annotation(s, 4.8, "从需求到全自动管道只用了一个对话回合。总裁提目标，AI团队完成从设计到部署的全链路。")

    # ═══════ SLIDE 7: 场景四 — 情报变行动 ═══════
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_gold_bar(s)
    add_label(s, "场景四 · 评审评分实录")
    add_title(s, "报告不是用来看的，是用来执行的", size=24)
    add_chat(s, 1.0, "总裁 {{PRESIDENT_NAME}}", "报告中的建议需要专业评估，批准后需要有执行者完成。我希望整个工作由多个专业AI Agent完成。", is_user=True)
    add_text(s, 0.6, 1.9, 8.8, 0.3, "于是建立了情报行动管线 — 每条建议经过4专家评审后才能执行：", font_size=12, color=MUTED)
    # Table data
    table_data = [
        ["建议", "战略", "决策", "趋势", "技术", "均分", "决定"],
        ["CLAUDE.md→Harness Config", "5", "4", "5", "5", "4.75", "批准"],
        ["创建方法论文档", "5", "4", "5", "—", "4.67", "批准"],
        ["QA自动化评分引擎", "4", "4", "5", "4", "4.25", "批准"],
        ["NVIDIA OpenShell", "—", "—", "—", "—", "—", "暂缓"],
    ]
    tbl = s.shapes.add_table(len(table_data), len(table_data[0]), Inches(0.6), Inches(2.25), Inches(8.8), Inches(1.6)).table
    for r, row in enumerate(table_data):
        for c, cell_text in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = "Calibri"
                p.font.color.rgb = WHITE if r == 0 else TEXT
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if r == 0 else CARD
    add_annotation(s, 4.1, "评审规则：均分≥4.0批准 / 3.0-3.9有条件 / <3.0暂缓 / 任一专家=1一票否决。3条批准后由Harness Ops执行。")

    # ═══════ SLIDE 8: 场景五 — 发现Bug ═══════
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_gold_bar(s)
    add_label(s, "场景五 · 在实践中发现漏洞")
    add_title(s, "总裁追问：刚才谁执行的？", size=24)
    add_text(s, 0.6, 1.1, 8.8, 0.3, "体系运行后，总裁的一句追问揭开了深层问题：", font_size=12, color=MUTED)
    add_chat(s, 1.45, "总裁 {{PRESIDENT_NAME}}", "刚才谁执行的？是AI团队吗？", is_user=True)
    add_chat(s, 2.25, "Lysander AI", "直说：不是。是我一个人执行的，没有体现团队分工。这和之前指出的问题是同一个。", is_user=False)
    add_chat(s, 3.1, "总裁 {{PRESIDENT_NAME}}", "不要光是期望，要强制按照流程来干活。要充分分析原因，不要光是期望。", is_user=True)
    add_text(s, 0.6, 4.0, 8.8, 0.3, "这个问题反复出现了3次 — 不是态度问题，是机制问题。", font_size=13, color=RED, bold=True)
    add_annotation(s, 4.4, "体系不靠下次注意，靠制度强制。这是整场对话中最重要的认知转折。")

    # ═══════ SLIDE 9: Harness 方法论验证 ═══════
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_gold_bar(s)
    add_label(s, "场景五 · Harness Engineering 的实战验证")
    add_title(s, "Agent = Model + Harness", size=24)
    add_text(s, 0.6, 1.1, 8.8, 0.35, "这恰好验证了 Harness Engineering（2026年AI工程新范式）的核心概念：", font_size=12, color=MUTED)
    # Left card
    add_rect(s, 0.6, 1.55, 4.2, 1.8, CARD)
    add_rect(s, 0.6, 1.55, 0.04, 1.8, RED)
    add_text(s, 0.8, 1.6, 3.5, 0.3, "只有 Guides（前馈控制）", font_size=12, color=RED, bold=True)
    add_text(s, 0.8, 1.9, 3.8, 1.3, "CLAUDE.md 写了路由到团队\n但没有机制检查是否真的路由了\n\n→ 纸上的规则 ≠ 执行的保证\n\nLysander 可以无阻碍地\n跳过团队派单直接干活", font_size=11, color=TEXT)
    # Right card
    add_rect(s, 5.2, 1.55, 4.2, 1.8, CARD)
    add_rect(s, 5.2, 1.55, 0.04, 1.8, GREEN)
    add_text(s, 5.4, 1.6, 3.5, 0.3, "Guide + Sensor 双保险", font_size=12, color=GREEN, bold=True)
    add_text(s, 5.4, 1.9, 3.8, 1.3, "Guide：强制团队派单表\n（执行前必须输出谁做什么）\n\nSensor：执行审计师前置检查\n（没有派单就阻止执行）\n\n与开场问候同级，不可省略", font_size=11, color=TEXT)
    # Bottom insight
    add_rect(s, 0.6, 3.6, 8.8, 0.55, CARD)
    add_rect(s, 0.6, 3.6, 0.04, 0.55, GOLD)
    add_text(s, 0.85, 3.6, 8.3, 0.55, "Harness Engineering: 模型不变，改变 Harness（环境、约束、反馈），Agent 的行为就改变。\n竞争力不在模型，在 Harness。 —— Martin Fowler / Red Hat / Anthropic, 2026", font_size=11, color=WHITE)

    # ═══════ SLIDE 10: 插曲 — 视觉升级 ═══════
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_gold_bar(s)
    add_label(s, "插曲 · 一个有趣的小案例")
    add_title(s, "参考一下官网 — AI逆向工程品牌色彩", size=22)
    add_chat(s, 1.0, "总裁 {{PRESIDENT_NAME}}", "视觉效果参考 https://{{YOUR_WEBSITE}}/ 重新定义该内容", is_user=True)
    add_chat(s, 1.8, "Lysander AI", "网站是动态渲染的，CSS无法直接提取。换一个方式 — 抓取 Logo SVG 源码，从中逆向工程品牌色。", is_user=False)
    add_text(s, 0.6, 2.7, 8.8, 0.3, "AI 团队从一个 SVG 文件中提取出完整品牌色彩体系：", font_size=12, color=MUTED)
    # Color swatches
    colors = [("#FCAD2A", "Gold 品牌金"), ("#013A7D", "Deep Blue"), ("#028CDC", "Cyan 天蓝"), ("#0A1628", "Dark BG")]
    for i, (hex_color, name) in enumerate(colors):
        x = 0.6 + i * 2.3
        r, g, b = int(hex_color[1:3], 16), int(hex_color[3:5], 16), int(hex_color[5:7], 16)
        add_rect(s, x, 3.05, 0.4, 0.4, RGBColor(r, g, b))
        add_text(s, x + 0.5, 3.05, 1.5, 0.4, f"{hex_color}\n{name}", font_size=10, color=TEXT)
    add_text(s, 0.6, 3.6, 8.8, 0.35, "然后用这套色彩重写了报告模板 — 头部、表格、优先级标签、页脚全部对齐品牌视觉。", font_size=12, color=TEXT)
    add_annotation(s, 4.05, "给AI一个URL，它自己完成调研→提取→设计→实现→验证。这就是AI团队的工作方式。")

    # ═══════ SLIDE 11: HR 管理 ═══════
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_gold_bar(s)
    add_label(s, "场景六 · Agent 也需要绩效考核")
    add_title(s, "首次审计：Janus全员0分，Growth全员0分", size=22)
    add_chat(s, 1.0, "总裁 {{PRESIDENT_NAME}}", "需要有专门的HR支持创建团队，创建的成员能力要符合标准，并且随着发展不断提升能力。", is_user=True)
    # Left: Problems
    add_rect(s, 0.6, 1.9, 4.2, 1.2, CARD)
    add_rect(s, 0.6, 1.9, 0.04, 1.2, RED)
    add_text(s, 0.8, 1.92, 2, 0.25, "审计发现", font_size=11, color=RED, bold=True)
    add_text(s, 0.8, 2.15, 3.8, 0.85, "· Janus 6人全部 0 分（无元数据）\n· Growth 2人全部 0 分（Schema缺失）\n· 三种不同卡片格式混用\n· 大量C级能力描述", font_size=10, color=TEXT)
    # Right: Solution
    add_rect(s, 5.2, 1.9, 4.2, 1.2, CARD)
    add_rect(s, 5.2, 1.9, 0.04, 1.2, GREEN)
    add_text(s, 5.4, 1.92, 2, 0.25, "建立的制度", font_size=11, color=GREEN, bold=True)
    add_text(s, 5.4, 2.15, 3.8, 0.85, "· 强制Schema — 缺项不得上岗\n· A/B/C能力分级 — C级不合格\n· 入职审批 — 8项检查清单\n· 周度自动审计 — 每周一自动评分", font_size=10, color=TEXT)
    # Progress bars
    add_text(s, 0.6, 3.25, 2, 0.25, "四轮整治：", font_size=11, color=GOLD, bold=True)
    rounds = [("第一轮·发现问题", 64.1), ("第二轮·基础修复", 83.1), ("第三轮·提标90", 88.6), ("最终·全员达标", 93.8)]
    for i, (label, score) in enumerate(rounds):
        y = 3.55 + i * 0.32
        add_text(s, 0.6, y, 2.2, 0.25, label, font_size=9, color=MUTED)
        w = 5 * (score / 100)
        add_rect(s, 2.9, y + 0.07, w, 0.12, GOLD)
        add_rect(s, 2.9 + w, y + 0.07, 5 - w, 0.12, RGBColor(0x1E, 0x29, 0x3B))
        add_text(s, 8.1, y, 0.8, 0.25, str(score), font_size=11, color=GOLD, bold=True)

    # ═══════ SLIDE 12: 命名风暴 ═══════
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_gold_bar(s)
    add_label(s, "场景七 · 头脑风暴现场")
    add_title(s, "20个候选 → 5维评审 → Synapse", size=24)
    add_text(s, 0.6, 1.1, 8.8, 0.3, "为整套体系命名，Lysander 组织智囊团全员头脑风暴 ——", font_size=12, color=MUTED)
    add_expert(s, 1.5, "关联发现", "5条隐喻路径（Janus品牌/神经系统/操作系统/建筑/数字孪生），产出20个候选名")
    add_expert(s, 1.95, "趋势洞察师", "行业调研：避免 OS 后缀（太多人用）、AI 前缀（无辨识度）。好名字=一个词能讲故事")
    add_expert(s, 2.4, "决策顾问", "5维评审：品牌关联25% / 记忆度25% / 含义深度20% / 差异化15% / 扩展性15%")
    # Table
    name_data = [
        ["候选", "品牌", "记忆", "含义", "差异", "扩展", "总分"],
        ["Janus Cortex", "5", "4", "5", "4", "4", "4.50"],
        ["Synapse ★", "3", "5", "5", "4", "4", "4.20"],
        ["Nexus", "3", "5", "4", "3", "5", "4.00"],
    ]
    tbl = s.shapes.add_table(len(name_data), len(name_data[0]), Inches(0.6), Inches(3.0), Inches(8.8), Inches(1.1)).table
    for r, row in enumerate(name_data):
        for c, cell_text in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = "Calibri"
                p.font.color.rgb = WHITE if r == 0 else (GOLD if r == 2 else TEXT)
                p.font.bold = (r == 0 or r == 2)
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if r == 0 else CARD
    add_chat(s, 4.25, "总裁 {{PRESIDENT_NAME}}", "Synapse", is_user=True)
    add_text(s, 0.6, 5.05, 8.8, 0.25, "突触 — 神经元之间传递信号的关键节点。知识←突触→决策←突触→执行。一个词，讲完了整个体系。", font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)

    # ═══════ SLIDE 13: 全景 ═══════
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); add_gold_bar(s)
    add_label(s, "SYNAPSE 全景")
    add_title(s, "10 团队 · 46 Agent · 5 自动任务 · 24/7 运转", size=22)
    stats = [("10", "专业团队"), ("46", "AI Agent"), ("93.8", "审计均分"), ("0", "不合格"), ("5", "自动任务"), ("24/7", "运转")]
    for i, (num, label) in enumerate(stats):
        x = 0.3 + i * 1.6
        add_text(s, x, 1.1, 1.4, 0.5, num, font_size=28, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(s, x, 1.55, 1.4, 0.25, label, font_size=10, color=MUTED, alignment=PP_ALIGN.CENTER)
    # Architecture
    arch_lines = [
        ("Synapse 体系架构", GOLD, True, 13),
        ("├── 记忆层 ── Obsidian 第二大脑 (OBS)", TEXT, False, 11),
        ("├── 控制层 ── Harness Engineering (Guides + Sensors)", TEXT, False, 11),
        ("├── 执行层 ── 10 个团队 46 个 AI Agent", TEXT, False, 11),
        ("│    Graphify(6)·HR(2)·Harness Ops(4)·Butler(7)·RD(5)·OBS(4)·Content(4)·Growth(2)·Janus(6)·Stock(5)", MUTED, False, 9),
        ("├── 进化层 ── 情报闭环 (日报→评估→执行→报告)", TEXT, False, 11),
        ("└── 决策层 ── L1自动 → L2专家 → L3 CEO → L4总裁", TEXT, False, 11),
    ]
    add_multiline(s, 0.6, 2.0, 8.8, 2.5, arch_lines, font_size=11)
    add_annotation(s, 4.7, "每天自动运转：6am任务恢复→8am情报日报→10am行动管线→每周一HR审计。总裁只看Slack通知。")

    # ═══════ SLIDE 14: 核心洞察 ═══════
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_label(s, "核心洞察")
    add_text(s, 0.6, 1.0, 8.8, 2.0, "AI 的竞争力不在模型本身，\n而在于你为它构建的\n环境、约束、反馈和组织。", font_size=30, color=WHITE, italic=True)
    add_rect(s, 0.6, 3.2, 2, 0.015, GOLD)
    add_text(s, 0.6, 3.4, 8, 0.4, "模型人人可用。Harness 才是壁垒。", font_size=18, color=GOLD, bold=True)
    add_text(s, 0.6, 3.9, 8, 0.4, "这不是理论 — 是我们一天内亲手搭建并反复修正的真实经历。", font_size=12, color=DIM)

    # ═══════ SLIDE 15: 结尾 ═══════
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_rect(s, 0, 5.3, 10, 0.04, GOLD)
    add_text(s, 0.6, 1.3, 8.8, 0.8, "SYNAPSE", font_size=44, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, 0.6, 2.2, 8.8, 0.4, "{{COMPANY_NAME}} AI Operating Framework", font_size=16, color=MUTED, alignment=PP_ALIGN.CENTER)
    add_rect(s, 4, 2.8, 2, 0.015, GOLD)
    add_text(s, 0.6, 3.1, 8.8, 0.3, "Powered by Claude Code · Harness Engineering · Obsidian", font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)
    add_text(s, 0.6, 3.6, 8.8, 0.3, "{{PRESIDENT_NAME}}  |  {{COMPANY_NAME}}  |  2026", font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)

    # ── Save ──
    out_path = Path(__file__).parent.parent / "obs" / "generated-articles" / "2026-04-10-synapse-presentation.pptx"
    prs.save(str(out_path))
    print(f"Generated: {out_path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Size: {out_path.stat().st_size // 1024} KB")


if __name__ == "__main__":
    main()
